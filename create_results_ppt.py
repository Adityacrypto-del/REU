#!/usr/bin/env python3
"""
Generate a comprehensive research presentation for the Neural Collapse REU project.
Includes actual experimental results and plot images from 50-epoch training.
"""

import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
#  Color Palette & Styling Constants
# ============================================================
BG_DARK      = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
BG_CARD      = RGBColor(0x1A, 0x25, 0x3C)   # Card background
ACCENT_BLUE  = RGBColor(0x3B, 0x82, 0xF6)   # Primary accent
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # Success green
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # Warning amber
ACCENT_RED   = RGBColor(0xEF, 0x44, 0x44)   # Alert red
TEXT_WHITE    = RGBColor(0xF1, 0xF5, 0xF9)   # Primary text
TEXT_GRAY     = RGBColor(0x94, 0xA3, 0xB8)   # Secondary text
TEXT_DIM      = RGBColor(0x64, 0x74, 0x8B)   # Dim text
LAYER_COLORS  = [
    RGBColor(0xE7, 0x4C, 0x3C),  # Layer 1 — red
    RGBColor(0xE6, 0x7E, 0x22),  # Layer 2 — orange
    RGBColor(0x2E, 0xCC, 0x71),  # Layer 3 — green
    RGBColor(0x34, 0x98, 0xDB),  # Layer 4 — blue
]

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a styled text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=TEXT_WHITE, spacing=Pt(6)):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_card(slide, left, top, width, height, color=BG_CARD):
    """Add a rounded card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
#  Slide Builders
# ============================================================

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    add_accent_line(slide, Inches(1), Inches(1.8), Inches(2), ACCENT_BLUE)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Investigating Neural Collapse Geometry &\nIts Implications for Early Exiting in CNNs",
                 font_size=32, bold=True, color=TEXT_WHITE)

    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                 "on Fine-Grained Visual Classification (Oxford Flowers 102)",
                 font_size=22, color=ACCENT_BLUE)

    add_accent_line(slide, Inches(1), Inches(4.3), Inches(1.5), ACCENT_GREEN)

    add_text_box(slide, Inches(1), Inches(4.6), Inches(5), Inches(0.4),
                 "Aditya Arasamangalam", font_size=20, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1), Inches(5.1), Inches(5), Inches(0.4),
                 "Research Experience for Undergraduates (REU)", font_size=16, color=TEXT_GRAY)
    add_text_box(slide, Inches(1), Inches(5.5), Inches(5), Inches(0.4),
                 "Guide: Sunil Gurlahosur", font_size=16, color=TEXT_GRAY)


def slide_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "MOTIVATION", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    cards = [
        ("🔍  The Black Box Problem",
         "CNNs achieve high accuracy but their intermediate representations remain opaque. "
         "Understanding layer-wise geometry is crucial for trust and deployment.",
         ACCENT_BLUE),
        ("⚡  Efficiency via Early Exiting",
         "If intermediate layers develop strong geometric structure, we can halt computation early — "
         "saving energy and latency for edge deployment and real-time inference.",
         ACCENT_GREEN),
        ("📊  Neural Collapse as a Lens",
         "NC describes a beautifully structured terminal state (NC1–NC4). All existing work studies it "
         "at the FINAL layer. Nobody has systematically tracked it through intermediate layers.",
         ACCENT_AMBER),
    ]

    for i, (title, desc, accent) in enumerate(cards):
        y = Inches(1.4) + Inches(i * 1.9)
        card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.6))
        add_accent_line(slide, Inches(0.8), y, Pt(4), accent)
        add_text_box(slide, Inches(1.2), y + Inches(0.15), Inches(10.5), Inches(0.4),
                     title, font_size=20, bold=True, color=TEXT_WHITE)
        add_text_box(slide, Inches(1.2), y + Inches(0.6), Inches(10.5), Inches(0.8),
                     desc, font_size=15, color=TEXT_GRAY)


def slide_problem_statement(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "PROBLEM STATEMENT", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    questions = [
        ("RQ1", "How does Neural Collapse evolve layer-by-layer through a ResNet-18 trained on Oxford Flowers 102?"),
        ("RQ2", "How does the strength of collapse at intermediate layers affect the reliability and accuracy of Early Exit classifiers?"),
        ("RQ3", "What is the interplay between geometric complexity, class imbalance, and the completeness of Neural Collapse?"),
    ]

    for i, (label, question) in enumerate(questions):
        y = Inches(1.4) + Inches(i * 1.8)
        card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.5))

        # RQ label badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.1), y + Inches(0.25), Inches(0.8), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = label
        badge_tf.paragraphs[0].font.size = Pt(14)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(slide, Inches(2.2), y + Inches(0.2), Inches(9.5), Inches(1.1),
                     question, font_size=17, color=TEXT_WHITE)


def slide_literature_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "LITERATURE REVIEW", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    papers = [
        ("Liu & Qin (CVPR 2025)", "OOD Detection via Neural Collapse",
         "Strong collapse geometry → reliable distance-based OOD detection",
         "Only studies final layer", ACCENT_BLUE),
        ("Wang et al. (CVPR 2024)", "Debiased Learning via NC",
         "Shortcut learning distorts ideal simplex geometry",
         "No intermediate layer bias analysis", ACCENT_GREEN),
        ("Munn et al. (arXiv 2024)", "Geometric Complexity in Transfer Learning",
         "Collapse strength depends on dataset geometric complexity",
         "No connection to early exiting", ACCENT_AMBER),
        ("Hasegawa & Sato (arXiv 2024)", "Multiplicative Logit Adjustment",
         "Class imbalance distorts collapse; logit scaling repairs it",
         "Only at the final classifier", ACCENT_RED),
    ]

    for i, (authors, title, finding, gap, accent) in enumerate(papers):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.1)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_card(slide, x, y, Inches(5.8), Inches(2.6))
        add_accent_line(slide, x, y, Inches(5.8), accent)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.3), Inches(0.35),
                     authors, font_size=13, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.3), Inches(0.35),
                     title, font_size=16, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(5.3), Inches(0.7),
                     f"✓  {finding}", font_size=13, color=TEXT_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.7), Inches(5.3), Inches(0.7),
                     f"⚠  Gap: {gap}", font_size=13, color=ACCENT_AMBER)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "MODEL ARCHITECTURE", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "EarlyExitResNet-18: ResNet-18 backbone with 4 early exit classifiers",
                 font_size=17, color=TEXT_GRAY)

    # Architecture flow: 4 layer blocks with exit heads
    layers = [
        ("Layer 1", "64-d", "56×56", LAYER_COLORS[0]),
        ("Layer 2", "128-d", "28×28", LAYER_COLORS[1]),
        ("Layer 3", "256-d", "14×14", LAYER_COLORS[2]),
        ("Layer 4", "512-d", "7×7", LAYER_COLORS[3]),
    ]

    for i, (name, dim, spatial, color) in enumerate(layers):
        x = Inches(0.8) + i * Inches(3.0)
        y = Inches(2.0)

        # Main block
        block = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.5), Inches(1.4)
        )
        block.fill.solid()
        block.fill.fore_color.rgb = BG_CARD
        block.line.color.rgb = color
        block.line.width = Pt(2)

        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(2.1), Inches(0.4),
                     name, font_size=18, bold=True, color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(2.1), Inches(0.3),
                     f"Features: {dim}", font_size=13, color=TEXT_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.1), Inches(0.3),
                     f"Spatial: {spatial}", font_size=13, color=TEXT_DIM)

        # Exit arrow + classifier
        exit_y = Inches(3.7)
        exit_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), exit_y, Inches(1.9), Inches(0.9)
        )
        exit_box.fill.solid()
        exit_box.fill.fore_color.rgb = color
        exit_box.line.fill.background()

        exit_tf = exit_box.text_frame
        exit_tf.word_wrap = True
        p = exit_tf.paragraphs[0]
        p.text = f"Exit {i+1}\nAdaptivePool → FC(102)"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Why ResNet section
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
                 "Why ResNet-18?", font_size=20, bold=True, color=ACCENT_GREEN)

    reasons = [
        "✓ Canonical architecture for Neural Collapse literature (directly comparable)",
        "✓ Clean residual blocks → natural layer-wise hierarchy for early exits",
        "✓ Skip connections preserve gradient flow → meaningful features even at early layers",
        "✓ Computationally efficient (11.7M params) → rapid experimentation on Apple Silicon MPS"
    ]
    add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(2.0),
                    reasons, font_size=13, color=TEXT_GRAY)


def slide_nc_metrics_explained(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "NEURAL COLLAPSE METRICS (NC1–NC4)", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    metrics = [
        ("NC1", "Within-Class Collapse", "Tr(Σ_W) / Tr(Σ_B)", "→ 0",
         "Features collapse to class means", ACCENT_RED),
        ("NC2", "Simplex ETF", "cos(μ_c, μ_c')", "→ −1/(C−1)",
         "Class means form equiangular tight frame", ACCENT_AMBER),
        ("NC3", "Classifier Alignment", "cos(W_c, μ_c − μ_G)", "→ 1.0",
         "Classifier weights align with class means", ACCENT_GREEN),
        ("NC4", "NCC Classification", "argmin ||h − μ_c||", "→ 100%",
         "Simple distance classifier matches trained classifier", ACCENT_BLUE),
    ]

    for i, (label, name, formula, target, desc, accent) in enumerate(metrics):
        col = i % 2
        row = i // 2
        x = Inches(0.8) + col * Inches(6.1)
        y = Inches(1.3) + row * Inches(2.9)

        card = add_card(slide, x, y, Inches(5.8), Inches(2.6))
        add_accent_line(slide, x, y, Inches(5.8), accent)

        # NC label badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = label
        badge_tf.paragraphs[0].font.size = Pt(14)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(slide, x + Inches(1.1), y + Inches(0.2), Inches(4.5), Inches(0.35),
                     name, font_size=18, color=TEXT_WHITE, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(5.3), Inches(0.4),
                     f"Formula: {formula}", font_size=14, color=TEXT_GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(5.3), Inches(0.4),
                     f"Target: {target}", font_size=14, color=accent, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.7), Inches(5.3), Inches(0.6),
                     desc, font_size=13, color=TEXT_DIM)


def slide_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "DATASET: OXFORD FLOWERS 102", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    props = [
        ("102", "Fine-Grained Classes", "Many flower species share similar petal shapes, colors, textures — "
         "creating high geometric complexity (per Munn et al.)", ACCENT_BLUE),
        ("~8,189", "Total Images", "Moderate scale — quality of learned geometry genuinely matters. "
         "Imbalance effects observable without over-parameterization masking.", ACCENT_GREEN),
        ("Rich", "Visual Hierarchy", "Strong low-level features (edges, colors) in early layers → "
         "increasingly abstract semantic features (species identity) in deeper layers.", ACCENT_AMBER),
    ]

    for i, (stat, label, desc, accent) in enumerate(props):
        y = Inches(1.3) + i * Inches(1.9)
        card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.6))

        # Large stat number
        add_text_box(slide, Inches(1.2), y + Inches(0.15), Inches(1.5), Inches(0.6),
                     stat, font_size=36, bold=True, color=accent)
        add_text_box(slide, Inches(1.2), y + Inches(0.7), Inches(1.5), Inches(0.3),
                     label, font_size=12, color=TEXT_DIM, bold=True)
        add_text_box(slide, Inches(3.0), y + Inches(0.2), Inches(9.0), Inches(1.2),
                     desc, font_size=15, color=TEXT_GRAY)


def slide_training_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "TRAINING CONFIGURATION", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    # Config table
    configs = [
        ("Architecture", "ResNet-18 (pretrained ImageNet) + 4 Early Exit heads"),
        ("Dataset", "Oxford Flowers 102 (train split → validation on test split)"),
        ("Optimizer", "Adam (lr = 1e-3)"),
        ("Scheduler", "CosineAnnealingLR (T_max = 50 per phase)"),
        ("Loss", "Cross-Entropy (sum of all 4 exits)"),
        ("Batch Size", "32"),
        ("Epochs", "50 (Phase 1) + 50 (Phase 2 — extended) = 100 total"),
        ("Device", "Apple Silicon M5 (MPS backend)"),
        ("NC Metrics", "Computed every epoch on validation set"),
        ("Checkpoints", "Saved every epoch (model + optimizer + scheduler + history)"),
    ]

    for i, (key, val) in enumerate(configs):
        y = Inches(1.3) + i * Inches(0.55)
        bg_color = BG_CARD if i % 2 == 0 else BG_DARK
        row = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(0.5)
        )
        row.fill.solid()
        row.fill.fore_color.rgb = bg_color
        row.line.fill.background()

        add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(2.5), Inches(0.4),
                     key, font_size=14, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, Inches(3.8), y + Inches(0.05), Inches(8.3), Inches(0.4),
                     val, font_size=14, color=TEXT_WHITE)


def slide_results_accuracy(prs, history):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "RESULTS: VALIDATION ACCURACY", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    # Best accuracy cards
    best_accs = []
    for i in range(4):
        best = max(history['val_accs'][e][i] for e in range(len(history['val_accs'])))
        best_epoch = max(range(len(history['val_accs'])), key=lambda e: history['val_accs'][e][i]) + 1
        best_accs.append((best, best_epoch))

    labels = ["Exit 1\n(Layer 1, 64-d)", "Exit 2\n(Layer 2, 128-d)",
              "Exit 3\n(Layer 3, 256-d)", "Exit 4\n(Layer 4, 512-d)"]

    for i in range(4):
        x = Inches(0.8) + i * Inches(3.0)
        y = Inches(1.3)
        card = add_card(slide, x, y, Inches(2.7), Inches(2.2))
        add_accent_line(slide, x, y, Inches(2.7), LAYER_COLORS[i])

        add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.4), Inches(0.7),
                     labels[i], font_size=13, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.4), Inches(0.7),
                     f"{best_accs[i][0]:.1%}", font_size=36, bold=True,
                     color=LAYER_COLORS[i], alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.6), Inches(2.4), Inches(0.4),
                     f"Best @ epoch {best_accs[i][1]}", font_size=12,
                     color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # Embed accuracy curves plot
    plot_path = './src/results/plots/accuracy_curves.png'
    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(0.8), Inches(3.8),
                                 width=Inches(11.5), height=Inches(3.4))


def slide_results_nc_evolution(prs, history):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "RESULTS: LAYER-WISE NEURAL COLLAPSE EVOLUTION",
                 font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    # Embed the NC evolution plot (the CENTRAL figure)
    plot_path = './src/results/plots/nc_evolution.png'
    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.2),
                                 width=Inches(12.3), height=Inches(6.0))


def slide_results_nc_table(prs, history):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "RESULTS: FINAL NC METRICS (EPOCH 50)", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    last_nc = history['nc_metrics'][-1]

    # Table header
    headers = ["Metric", "Layer 1 (64-d)", "Layer 2 (128-d)", "Layer 3 (256-d)", "Layer 4 (512-d)", "Trend"]
    col_widths = [Inches(2.5), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)]
    x_start = Inches(0.8)
    y_header = Inches(1.4)

    # Header row
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x_start, y_header, Inches(11.5), Inches(0.5)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = ACCENT_BLUE
    header_bg.line.fill.background()

    x = x_start
    for i, (header, w) in enumerate(zip(headers, col_widths)):
        add_text_box(slide, x + Inches(0.1), y_header + Inches(0.05), w, Inches(0.4),
                     header, font_size=13, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        x += w

    # Data rows
    rows_data = [
        ("NC1 (Sw/Sb) ↓ better",
         [f"{last_nc['layers'][l]['nc1']:.3f}" for l in range(4)],
         "✅ Monotonic ↓"),
        ("NC2 (cos sim → −0.0099)",
         [f"{last_nc['layers'][l]['nc2_mean']:.4f}" for l in range(4)],
         "✅ Approaching"),
        ("NC3 (alignment → 1.0)",
         [f"{last_nc['layers'][l]['nc3']:.3f}" for l in range(4)],
         "✅ Strong"),
        ("NC4 (NCC acc ↑ better)",
         [f"{last_nc['layers'][l]['nc4']:.1%}" for l in range(4)],
         "✅ Monotonic ↑"),
    ]

    for r, (metric_name, values, trend) in enumerate(rows_data):
        y = y_header + Inches(0.55) + r * Inches(0.6)
        bg_color = BG_CARD if r % 2 == 0 else BG_DARK
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x_start, y, Inches(11.5), Inches(0.55)
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = bg_color
        row_bg.line.fill.background()

        x = x_start
        add_text_box(slide, x + Inches(0.1), y + Inches(0.08), col_widths[0], Inches(0.4),
                     metric_name, font_size=13, color=TEXT_WHITE, bold=True)
        x += col_widths[0]

        for j, val in enumerate(values):
            add_text_box(slide, x + Inches(0.1), y + Inches(0.08), col_widths[j+1], Inches(0.4),
                         val, font_size=14, color=LAYER_COLORS[j], alignment=PP_ALIGN.CENTER)
            x += col_widths[j+1]

        add_text_box(slide, x + Inches(0.1), y + Inches(0.08), col_widths[5], Inches(0.4),
                     trend, font_size=12, color=ACCENT_GREEN, alignment=PP_ALIGN.CENTER)

    # Key insight
    add_card(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.5), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(4.4), Inches(10.5), Inches(0.4),
                 "🔑  KEY FINDING", font_size=18, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, Inches(1.2), Inches(4.9), Inches(10.5), Inches(0.7),
                 "Neural Collapse strengthens monotonically from shallow → deep layers. "
                 "NC1 decreases (0.783 → 0.584) and NC4 increases (29.5% → 85.5%) across all 4 residual blocks, "
                 "confirming the progressive collapse hypothesis on a fine-grained dataset.",
                 font_size=15, color=TEXT_GRAY)

    # NC vs Accuracy plot
    plot_path = './src/results/plots/nc_vs_accuracy.png'
    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(0.8), Inches(5.9),
                                 width=Inches(11.5), height=Inches(1.4))


def slide_results_early_exit(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "RESULTS: EARLY EXIT SIMULATION", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    # Threshold sweep table
    headers = ["Threshold", "Accuracy", "Avg Layers", "Speedup", "L1 Exits", "L2 Exits", "L3 Exits", "L4 Exits"]
    data_rows = [
        ("0.50", "82.03%", "3.20", "1.25×", "3", "506", "3,895", "1,745"),
        ("0.70", "82.58%", "3.41", "1.17×", "0", "196", "3,243", "2,710"),
        ("0.90", "82.65%", "3.66", "1.09×", "0", "19", "2,063", "4,067"),
        ("0.99", "82.62%", "3.92", "1.02×", "0", "0", "502", "5,647"),
    ]

    col_widths = [Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3),
                  Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3)]
    x_start = Inches(0.8)
    y_header = Inches(1.3)

    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x_start, y_header, Inches(11.5), Inches(0.5)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = ACCENT_BLUE
    header_bg.line.fill.background()

    x = x_start
    for header, w in zip(headers, col_widths):
        add_text_box(slide, x, y_header + Inches(0.05), w, Inches(0.4),
                     header, font_size=12, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        x += w

    for r, row in enumerate(data_rows):
        y = y_header + Inches(0.55) + r * Inches(0.5)
        bg = BG_CARD if r % 2 == 0 else BG_DARK
        rbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, y, Inches(11.5), Inches(0.45))
        rbg.fill.solid()
        rbg.fill.fore_color.rgb = bg
        rbg.line.fill.background()

        x = x_start
        for j, (val, w) in enumerate(zip(row, col_widths)):
            clr = ACCENT_GREEN if j == 3 else TEXT_WHITE
            add_text_box(slide, x, y + Inches(0.05), w, Inches(0.35),
                         val, font_size=13, color=clr, alignment=PP_ALIGN.CENTER)
            x += w

    # Insights
    add_card(slide, Inches(0.8), Inches(3.8), Inches(5.4), Inches(3.2), BG_CARD)
    add_text_box(slide, Inches(1.1), Inches(4.0), Inches(4.8), Inches(0.4),
                 "📊 Key Observations", font_size=18, bold=True, color=ACCENT_GREEN)
    insights = [
        "• At τ=0.50: Exit 3 handles 63% of all samples (3,895/6,149)",
        "• Accuracy stays stable (~82.6%) across all thresholds",
        "• Exit 1 almost never triggers — NC too weak at 64-d",
        "• Exit 3 is the \"sweet spot\" — strong NC, good accuracy",
        "• 1.25× speedup achievable with <1% accuracy loss",
    ]
    add_bullet_list(slide, Inches(1.1), Inches(4.5), Inches(4.8), Inches(2.5),
                    insights, font_size=13, color=TEXT_GRAY, spacing=Pt(4))

    # Entropy insight card
    add_card(slide, Inches(6.5), Inches(3.8), Inches(5.8), Inches(3.2), BG_CARD)
    add_text_box(slide, Inches(6.8), Inches(4.0), Inches(5.2), Inches(0.4),
                 "🔥 Shannon Entropy Analysis", font_size=18, bold=True, color=ACCENT_AMBER)
    entropy_data = [
        "Exit 1:  H = 4.06  (near-random — 102 classes)",
        "Exit 2:  H = 3.24  (still uncertain)",
        "Exit 3:  H = 1.14  (confident decisions)",
        "Exit 4:  H = 0.43  (highly confident)",
        "",
        "Entropy drops 10× from Layer 1 → Layer 4,",
        "directly mirroring NC strength progression."
    ]
    add_bullet_list(slide, Inches(6.8), Inches(4.5), Inches(5.2), Inches(2.5),
                    entropy_data, font_size=13, color=TEXT_GRAY, spacing=Pt(3))


def slide_results_ood(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
                 "RESULTS: OOD DETECTION (LIU & QIN 2025)",
                 font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "In-Distribution: Flowers-102  |  Out-of-Distribution: CIFAR-10",
                 font_size=16, color=TEXT_GRAY)

    # OOD results as visual cards
    ood_data = [
        ("Exit 1", "−0.081", "OOD closer than ID!", False, LAYER_COLORS[0]),
        ("Exit 2", "−0.047", "OOD still closer", False, LAYER_COLORS[1]),
        ("Exit 3", "+0.056", "ID closer — detection works!", True, LAYER_COLORS[2]),
        ("Exit 4", "+0.139", "Strong separation!", True, LAYER_COLORS[3]),
    ]

    for i, (exit_name, gap, interpretation, is_positive, color) in enumerate(ood_data):
        x = Inches(0.8) + i * Inches(3.0)
        y = Inches(1.9)
        card = add_card(slide, x, y, Inches(2.7), Inches(2.8))

        indicator_color = ACCENT_GREEN if is_positive else ACCENT_RED
        add_accent_line(slide, x, y, Inches(2.7), indicator_color)

        add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.4), Inches(0.4),
                     exit_name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.4), Inches(0.6),
                     f"Gap: {gap}", font_size=28, bold=True,
                     color=indicator_color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.4), Inches(2.4), Inches(0.4),
                     "✓" if is_positive else "✗", font_size=36,
                     color=indicator_color, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.9), Inches(2.4), Inches(0.7),
                     interpretation, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Conclusion box
    add_card(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0), BG_CARD)
    add_text_box(slide, Inches(1.2), Inches(5.2), Inches(10.5), Inches(0.4),
                 "🔑  OOD DETECTION VALIDATES PROGRESSIVE COLLAPSE",
                 font_size=18, bold=True, color=ACCENT_GREEN)
    add_text_box(slide, Inches(1.2), Inches(5.7), Inches(10.5), Inches(1.1),
                 "OOD detection ONLY works at Layers 3–4 (positive gap = ID data closer to class centers). "
                 "Layers 1–2 CANNOT distinguish in-distribution from out-of-distribution data. "
                 "This directly confirms that OOD detection reliability requires strong Neural Collapse geometry — "
                 "exactly as Liu & Qin (2025) theorized, but now validated at intermediate layers for the first time.",
                 font_size=15, color=TEXT_GRAY)


def slide_loss_curve(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "RESULTS: TRAINING & VALIDATION LOSS", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    plot_path = './src/results/plots/loss_curve.png'
    if os.path.exists(plot_path):
        slide.shapes.add_picture(plot_path, Inches(1.5), Inches(1.3),
                                 width=Inches(10), height=Inches(5.8))


def slide_key_findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "KEY FINDINGS", font_size=28, bold=True, color=ACCENT_GREEN)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_GREEN)

    findings = [
        ("1", "Progressive Collapse Confirmed",
         "NC1–NC4 metrics improve monotonically from Layer 1 → Layer 4, confirming that Neural Collapse develops progressively through intermediate residual blocks — not just at the final layer.",
         ACCENT_BLUE),
        ("2", "NC Strength Predicts Exit Reliability",
         "Strong positive correlation between NC4 (NCC accuracy) and exit classifier accuracy. Layers with NC4 > 50% produce usable early exits; below that, predictions are unreliable.",
         ACCENT_GREEN),
        ("3", "OOD Detection Requires Deep Collapse",
         "Distance-based OOD detection fails completely at Layers 1–2 (negative gap) but succeeds at Layers 3–4 (positive gap), providing the first empirical evidence of layer-dependent OOD reliability.",
         ACCENT_AMBER),
        ("4", "Exit 3 is the Practical Sweet Spot",
         "Layer 3 (256-d) achieves 77.3% accuracy with strong NC metrics, handles 63% of samples at τ=0.50, and provides meaningful OOD detection — making it the optimal early exit point.",
         ACCENT_RED),
    ]

    for i, (num, title, desc, accent) in enumerate(findings):
        y = Inches(1.3) + i * Inches(1.45)
        card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.3))
        add_accent_line(slide, Inches(0.8), y, Pt(4), accent)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.2), Inches(0.5), Inches(0.5)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = accent
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = num
        badge_tf.paragraphs[0].font.size = Pt(16)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = TEXT_WHITE
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(slide, Inches(1.9), y + Inches(0.1), Inches(10), Inches(0.35),
                     title, font_size=17, bold=True, color=TEXT_WHITE)
        add_text_box(slide, Inches(1.9), y + Inches(0.5), Inches(10), Inches(0.7),
                     desc, font_size=13, color=TEXT_GRAY)


def slide_future_work(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 "FUTURE WORK & NEXT STEPS", font_size=28, bold=True, color=ACCENT_BLUE)
    add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(1.5), ACCENT_BLUE)

    items = [
        ("🔄", "Extended Training (100 epochs)",
         "Currently running: Phase 2 training with fresh cosine LR schedule to push accuracy and NC metrics further.",
         ACCENT_GREEN),
        ("📊", "CIFAR-10 Comparative Study",
         "Run identical pipeline on CIFAR-10 (10 classes, low geometric complexity) to validate Munn et al.'s prediction that simpler datasets exhibit stronger collapse.",
         ACCENT_BLUE),
        ("⚖️", "Multiplicative Logit Adjustment at Early Exits",
         "Apply Hasegawa & Sato's MLA to intermediate classifiers — does it help underrepresented classes at shallow layers?",
         ACCENT_AMBER),
        ("🔬", "t-SNE / PCA Feature Visualization",
         "Visualize feature clustering at each layer to qualitatively complement the NC1–NC4 quantitative analysis.",
         ACCENT_RED),
        ("🏗️", "Architecture Extension",
         "Extend to DenseNet (feature reuse effects), EfficientNet (scaling effects), and ViTs (attention vs convolution).",
         TEXT_GRAY),
    ]

    for i, (icon, title, desc, accent) in enumerate(items):
        y = Inches(1.3) + i * Inches(1.15)
        card = add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.0))

        add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(0.5), Inches(0.4),
                     icon, font_size=20, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(1.7), y + Inches(0.05), Inches(10.2), Inches(0.35),
                     title, font_size=16, bold=True, color=accent)
        add_text_box(slide, Inches(1.7), y + Inches(0.4), Inches(10.2), Inches(0.55),
                     desc, font_size=12, color=TEXT_GRAY)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
                 "Thank You", font_size=48, bold=True, color=TEXT_WHITE,
                 alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5.5), Inches(3.8), Inches(2.3), ACCENT_BLUE)

    add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
                 "Aditya Arasamangalam  •  REU  •  Guide: Sunil Gurlahosur",
                 font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
                 "Questions?",
                 font_size=24, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER, bold=True)


# ============================================================
#  Main Builder
# ============================================================

def main():
    # Load training history
    history_path = './src/results/metrics/full_history.json'
    with open(history_path) as f:
        history = json.load(f)
    print(f"Loaded history: {len(history['train_loss'])} epochs")

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    print("Building slides...")
    slide_title(prs)                          # 1
    slide_motivation(prs)                     # 2
    slide_problem_statement(prs)              # 3
    slide_literature_overview(prs)            # 4
    slide_nc_metrics_explained(prs)           # 5
    slide_dataset(prs)                        # 6
    slide_architecture(prs)                   # 7
    slide_training_setup(prs)                 # 8
    slide_results_accuracy(prs, history)      # 9
    slide_loss_curve(prs)                     # 10
    slide_results_nc_evolution(prs, history)  # 11
    slide_results_nc_table(prs, history)      # 12
    slide_results_early_exit(prs)             # 13
    slide_results_ood(prs)                    # 14
    slide_key_findings(prs)                   # 15
    slide_future_work(prs)                    # 16
    slide_thank_you(prs)                      # 17

    out_path = './Neural_Collapse_REU_Presentation.pptx'
    prs.save(out_path)
    print(f"\n✅ Presentation saved to: {out_path}")
    print(f"   Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
