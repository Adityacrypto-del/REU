import os
import re
import collections
from pptx import Presentation
from pptx.util import Inches, Pt
import base64
import requests

markdown_content = """
## SLIDE 1: Title
**Title:** Layer-wise Analysis of Feature Evolution in Convolutional Neural Networks
**Subtitle:** Understanding Feature Representation and Neural Collapse
**Author:** Aditya Arasamangalam
**Course:** ReU
**Guide:** Sunil Gurlahosur

## SLIDE 2: Motivation
*   **The Black Box Problem:** CNNs achieve high accuracy but lack internal interpretability.
*   **Output-Centric Bias:** Current research predominantly evaluates final output classifications, ignoring intermediate feature structures.
*   **Trust & Explainability:** Analyzing how features evolve across hidden layers builds reliable, transparent models for edge deployment.

## SLIDE 3: Problem Statement & Research Question
*   **Research Gap:** A lack of quantifiable understanding regarding layer-by-layer feature behavior and how dataset geometric complexity influences learning dynamics.
*   **Core Question:** How do internal features geometrically evolve across convolution layers, and how do dataset properties accelerate or hinder this evolution?

## SLIDE 4: Research Objectives
1.  **Analyze** feature maps across intermediate convolution layers.
2.  **Quantify** progression using Neural Collapse metrics (NC1-NC4).
3.  **Evaluate** Early Exit reliability based on intermediate geometric structure.
4.  **Compare** behavior systematically across datasets with varying structures.

## SLIDE 5: Literature 1 — OOD Detection via Neural Collapse
**Paper:** Liu & Qin, 2025
*   **Concept:** Terminal-phase deep features form tightly packed clusters. This geometry acts as a strict boundary for Out-of-Distribution (OOD) detection.
*   **Limitation Addressed:** Focuses only on final layers; our work extends this to early exits.

**Methodology Diagram:**
```mermaid
flowchart LR
    A[Input Sample] --> B[CNN Feature Extraction]
    B --> C[Compute Distance to NC Center]
    C --> D{Distance > \n Threshold}
    D -- Yes --> E[Flag as OOD]
    D -- No --> F[Classify In-Distribution]
```

## SLIDE 6: Literature 2 — Debiased Learning
**Paper:** Wang et al., 2024
*   **Concept:** Shortcut learning distorts the ideal simplex geometry of Neural Collapse. Explicit geometric regularization ensures fair, unbiased representations.
*   **Limitation Addressed:** Bias propagation through *intermediate* architectural blocks.

**Methodology Diagram:**
```mermaid
flowchart LR
    A[Biased Input Data] --> B[Intermediate Feature Maps]
    B --> C[Measure Geometric Skew]
    C --> D[Apply Orthogonal Regularization]
    D --> E[Unbiased Symmetric Features]
```

## SLIDE 7: Literature 3 — Geometric Complexity
**Paper:** Munn et al., 2024
*   **Concept:** Neural Collapse is not guaranteed. It is strictly dependent on the "Geometric Complexity" (target dataset's difficulty).
*   **Rationale:** Complex, fine-grained datasets yield weak representation overlaps, making them harder to collapse.

**Methodology Diagram:**
```mermaid
flowchart TD
    A[Target Dataset] --> B{Geometric Complexity?}
    B -- Low CIFAR-10 --> C[Strong Complete Collapse]
    B -- High Flowers 102 --> D[Weak Partial Collapse]
    C --> E[High Early Exit Reliability]
    D --> F[Low Early Exit Reliability]
```

## SLIDE 8: Literature 4 — Logit Adjustment
**Paper:** Hasegawa & Sato, 2024
*   **Concept:** Class imbalance mechanically shifts classification boundaries away from optimal Neural Collapse geometry. Multiplicative scaling can mathematically repair this.

**Methodology Diagram:**
```mermaid
flowchart LR
    A[Imbalanced Training Phase] --> B[Distorted Class Boundaries]
    B --> C[Compute Relative Weight Norms]
    C --> D[Apply Multiplicative Scaling]
    D --> E[Restored Symmetric Boundaries]
```

## SLIDE 9: Dataset Selection
*   **Primary Focus: Oxford Flowers 102**
    *   **Balance:** 102 classes with moderate data scale.
    *   **Complexity:** Fine-grained geometry (high inter-class similarity) makes it ideal for testing partial Neural Collapse.
*   **Comparative Baselines:**
    *   Food-101 (Texture-dominant)
    *   Stanford Cars (Geometry-dominant)
    *   FGVC Aircraft (Spatial/Structure-dominant)

## SLIDE 10: Model Architecture — ResNet-18
*   **Why ResNet?**
    *   **Residual Connections:** Prevents vanishing gradients, allowing stable measurement of deep features.
    *   **Distinct Hierarchical Blocks:** Perfect for strategically attaching intermediate Early Exit classifiers.
*   **Comparison:** Unlike DenseNet (feature entanglement) or VGG (no residuals), ResNet provides clean architectural benchmarks.

## SLIDE 11: Proposed 5-Stage Framework
**Methodology Diagram:**
```mermaid
flowchart TD
    A[1. Feature Extraction: Hook Intermediate Layers] --> B[2. Statistical Analysis: Layer-wise Mean & Variance]
    B --> C[3. Neural Collapse: Quantify NC1-NC4]
    C --> D[4. Cross-Dataset Study: Compare Architectures]
    D --> E[5. Insights & Validation: Early Exit Reliability]
```

## SLIDE 12: Expected Results
*   **Early Layers:** Responsive to raw textures and dense edges. Weak collapse limits prediction confidence.
*   **Mid Layers:** Capture gradients and shapes. Logit adjustments stabilize boundary distortions.
*   **Deep Layers:** Form structured semantic geometries. High collapse guarantees reliable, calibration-sound inferences.

## SLIDE 13: Conclusion & Next Steps
*   **Key Contribution:** Developing a principled framework linking internal convolutional geometry directly to inference cost-savings (Early Exiting).
*   **Future Work:**
    *   Detailed feature clustering analysis via t-SNE / PCA.
    *   Expansion to Vision Transformers (ViTs) to compare spatial architectures against self-attention collapse.
"""

def generate_mermaid_image(mermaid_str, filename):
    try:
        # Encode strictly based on mermaid.ink format spec
        mermaid_encoded = base64.urlsafe_b64encode(mermaid_str.encode('utf8')).decode('ascii')
        url = f"https://mermaid.ink/img/{mermaid_encoded}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            with open(filename, 'wb') as f:
                f.write(response.content)
            return True
        return False
    except Exception as e:
        print(f"Mermaid generation failed: {e}")
        return False

def parse_markdown(md):
    slides = []
    current_slide = None
    lines = md.strip().split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith("## SLIDE"):
            if current_slide is not None:
                slides.append(current_slide)
            title = line.split(":", 1)[1].strip() if ":" in line else line
            current_slide = {'title': title, 'bullets': [], 'mermaid': None}
        elif current_slide is not None:
            if line.startswith("```mermaid"):
                mermaid_str = ""
                i += 1
                while i < len(lines) and not lines[i].startswith("```"):
                    mermaid_str += lines[i] + "\n"
                    i += 1
                current_slide['mermaid'] = mermaid_str.strip()
            elif line:
                # Remove strong bold
                clean_line = line.replace("**", "").replace("*   ", "").replace("1.  ", "").replace("2.  ", "").replace("3.  ", "").replace("4.  ", "").replace("5.  ", "")
                if clean_line.strip().startswith("Methodology Diagram"):
                    pass # ignore
                elif clean_line:
                    current_slide['bullets'].append(clean_line.strip())
        i += 1
    
    if current_slide is not None:
        slides.append(current_slide)
    return slides

def make_ppt():
    template_path = r"C:\Users\91948\Downloads\PPT Template.pptx"
    out_path = r"C:\Users\91948\OneDrive\Attachments\ReU\Final_Presentation.pptx"
    
    try:
        prs = Presentation(template_path)
        print("Using provided PPT Template.")
    except Exception as e:
        print(f"Failed to open template because: {e}. Falling back to default Presentation.")
        prs = Presentation()
        
    slides_data = parse_markdown(markdown_content)
    
    for idx, slide_info in enumerate(slides_data):
        layout_idx = 0 if idx == 0 else min(1, len(prs.slide_layouts)-1)
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        
        shapes = slide.shapes
        title_shape = shapes.title
        if title_shape:
            title_shape.text = slide_info['title']
            
        # Add bullets
        if len(slide_info['bullets']) > 0:
            if len(shapes.placeholders) > 1:
                body_shape = shapes.placeholders[1]
            else:
                # Add text box if no placeholder
                body_shape = shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
                
            tf = body_shape.text_frame
            for i, bullet in enumerate(slide_info['bullets']):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet
                if len(slide_info['bullets']) > 5:
                    p.font.size = Pt(14)
                else:
                    p.font.size = Pt(18)

        # Add mermaid
        if slide_info['mermaid']:
            img_path = f"mermaid_{idx}.png"
            if generate_mermaid_image(slide_info['mermaid'], img_path):
                # Put image at the bottom or middle based on text existence
                top = Inches(3.5) if slide_info['bullets'] else Inches(1.5)
                left = Inches(1.5)
                slide.shapes.add_picture(img_path, left, top, width=Inches(7))
            else:
                print(f"Could not generate image for slide {idx}")
                
    prs.save(out_path)
    print(f"Successfully created presentation at {out_path}")

if __name__ == "__main__":
    make_ppt()
